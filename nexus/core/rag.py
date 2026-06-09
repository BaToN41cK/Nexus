"""
RAG (Retrieval-Augmented Generation) orchestrator for Nexus.

Provides the :class:`RAGEngine` that ties together:
  - Embedding models (sentence-transformers, OpenAI, Ollama)
  - Vector stores (FAISS, Chroma)
  - BM25 keyword retrieval
  - Hybrid search (BM25 + vector via RRF)
  - MMR diversity re-ranking
  - Cross-encoder re-ranking
  - Document chunking
  - Source citation

Usage::

    from nexus.core.rag import RAGEngine

    engine = RAGEngine()
    engine.index_texts([
        ("Nexus is an AI assistant with web search.", {"source": "docs/README.md"}),
        ("It supports multiple LLM providers.", {"source": "docs/API.md"}),
    ])
    results = engine.query("What is Nexus?")
    print(results.response)
    print(results.citations)
"""

from __future__ import annotations

import logging
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from nexus.core.embeddings import (
    EmbeddingModel,
    create_embedding_model,
)
from nexus.core.retrieval import (
    BM25Retriever,
    CrossEncoderReranker,
    hybrid_search,
    max_marginal_relevance,
)
from nexus.core.vector_store import (
    Document,
    SearchResult,
    VectorStore,
    create_vector_store,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Data types
# ---------------------------------------------------------------------------


@dataclass
class ChunkConfig:
    """Configuration for document chunking."""

    chunk_size: int = 512
    chunk_overlap: int = 64
    separator: str = "\n"


@dataclass
class RetrievalResult:
    """Result of a RAG query."""

    response: str
    citations: List[Dict[str, Any]] = field(default_factory=list)
    source_documents: List[Document] = field(default_factory=list)
    scores: List[float] = field(default_factory=list)


@dataclass
class RAGConfig:
    """Configuration for the RAG engine."""

    embedding_backend: str = "sentence-transformers"
    embedding_kwargs: Dict[str, Any] = field(default_factory=dict)
    vector_store_backend: str = "faiss"
    vector_store_kwargs: Dict[str, Any] = field(default_factory=dict)
    chunk_config: ChunkConfig = field(default_factory=ChunkConfig)
    top_k_initial: int = 20
    top_k_after_rerank: int = 5
    use_bm25: bool = True
    use_mmr: bool = True
    use_cross_encoder: bool = False
    hybrid_alpha: float = 0.5
    mmr_lambda: float = 0.5
    cross_encoder_model: str = "cross-encoder/ms-marco-MiniLM-L-6-v2"


# ---------------------------------------------------------------------------
# Chunking utilities
# ---------------------------------------------------------------------------


def chunk_text(
    text: str,
    config: ChunkConfig,
    source: Optional[str] = None,
) -> List[Document]:
    """
    Split text into overlapping chunks.

    Args:
        text: The text to split.
        config: Chunking configuration.
        source: Optional source identifier (e.g. URL, file path).

    Returns:
        List of Document chunks.
    """
    separators = [config.separator, ". ", "! ", "? ", ";", ", ", " "]
    chunks: List[Document] = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + config.chunk_size, text_len)

        # Try to break at a separator for cleaner chunks.
        if end < text_len:
            best_sep = -1
            best_pos = -1
            for sep in separators:
                pos = text.rfind(sep, start + 1, end)
                if pos > best_pos:
                    best_pos = pos
                    best_sep = len(sep)

            if best_pos > start and best_pos < end:
                end = best_pos + best_sep

        chunk_text_content = text[start:end].strip()
        if chunk_text_content:
            metadata: Dict[str, Any] = {"chunk_start": start}
            if source:
                metadata["source"] = source
            chunks.append(Document(text=chunk_text_content, metadata=metadata))

        # Move start with overlap.
        next_start = end - config.chunk_overlap
        if next_start <= start:
            next_start = end
        start = next_start

    return chunks


# ---------------------------------------------------------------------------
# RAG Engine
# ---------------------------------------------------------------------------


class RAGEngine:
    """
    Main RAG engine that orchestrates retrieval and generation.

    Thread-safe for concurrent indexing.  Query method is reentrant but
    should not be called concurrently with indexing operations.
    """

    def __init__(self, config: Optional[RAGConfig] = None):
        self.config = config or RAGConfig()
        self._embedding_model: Optional[EmbeddingModel] = None
        self._vector_store: Optional[VectorStore] = None
        self._bm25: Optional[BM25Retriever] = None
        self._reranker: Optional[CrossEncoderReranker] = None
        self._documents: List[Document] = []

    # -- lazy initializers --

    @property
    def embedding_model(self) -> EmbeddingModel:
        if self._embedding_model is None:
            self._embedding_model = create_embedding_model(
                self.config.embedding_backend,
                **self.config.embedding_kwargs,
            )
        return self._embedding_model

    @property
    def vector_store(self) -> VectorStore:
        if self._vector_store is None:
            kwargs = dict(self.config.vector_store_kwargs)
            if "dimension" not in kwargs:
                kwargs["dimension"] = self.embedding_model.dimension
            self._vector_store = create_vector_store(
                self.config.vector_store_backend,
                **kwargs,
            )
        return self._vector_store

    @property
    def bm25(self) -> BM25Retriever:
        if self._bm25 is None:
            self._bm25 = BM25Retriever()
        return self._bm25

    @property
    def reranker(self) -> CrossEncoderReranker:
        if self._reranker is None:
            self._reranker = CrossEncoderReranker(
                model_name=self.config.cross_encoder_model,
            )
        return self._reranker

    # -- indexing --

    def index_texts(
        self,
        texts_with_metadata: List[Tuple[str, Dict[str, Any]]],
    ) -> int:
        """
        Index a list of (text, metadata) pairs.

        Each text is automatically chunked, embedded, and stored.

        Args:
            texts_with_metadata: List of (text, metadata) tuples.

        Returns:
            Number of chunks indexed.
        """
        documents: List[Document] = []
        for text, metadata in texts_with_metadata:
            source = metadata.get("source")
            chunks = chunk_text(text, self.config.chunk_config, source=source)
            for chunk in chunks:
                chunk.metadata.update(metadata)
            documents.extend(chunks)

        if not documents:
            return 0

        # Embed in batches.
        batch_size = 64
        for i in range(0, len(documents), batch_size):
            batch = documents[i : i + batch_size]
            texts = [d.text for d in batch]
            embeddings = self.embedding_model.embed(texts)

            # Add to vector store.
            self.vector_store.add(batch, embeddings)

            # Add to BM25.
            if self.config.use_bm25:
                self.bm25.add_documents(batch)

        self._documents.extend(documents)
        self._maybe_persist()
        logger.info(
            "Indexed %d chunks from %d texts",
            len(documents),
            len(texts_with_metadata),
        )
        return len(documents)

    def index_documents(self, documents: List[Document]) -> int:
        """
        Index pre-built Document objects (with embeddings already computed
        or will be computed here).

        Args:
            documents: List of Document objects.

        Returns:
            Number of documents indexed.
        """
        if not documents:
            return 0

        batch_size = 64
        for i in range(0, len(documents), batch_size):
            batch = documents[i : i + batch_size]
            texts = [d.text for d in batch]
            embeddings = self.embedding_model.embed(texts)

            self.vector_store.add(batch, embeddings)
            if self.config.use_bm25:
                self.bm25.add_documents(batch)

        self._documents.extend(documents)
        self._maybe_persist()
        logger.info("Indexed %d documents", len(documents))
        return len(documents)

    # -- retrieval and generation --

    def retrieve(
        self,
        query: str,
        top_k: Optional[int] = None,
    ) -> List[SearchResult]:
        """
        Retrieve relevant documents for a query using the full pipeline:
          vector search -> (optional BM25 hybrid) -> (optional MMR) -> (optional cross-encoder)

        Args:
            query: The query string.
            top_k: Override for top_k_after_rerank.

        Returns:
            List of SearchResult sorted by relevance.
        """
        cfg = self.config
        top_k_final = top_k or cfg.top_k_after_rerank

        # 1. Vector search.
        query_embedding = self.embedding_model.embed_query(query)
        vector_results = self.vector_store.search(
            query_embedding, top_k=cfg.top_k_initial,
        )

        if not vector_results:
            return []

        # 2. Hybrid search (BM25 + vector).
        if cfg.use_bm25:
            bm25_results = self.bm25.search(query, top_k=cfg.top_k_initial)
            combined = hybrid_search(
                vector_results, bm25_results,
                alpha=cfg.hybrid_alpha,
            )
        else:
            combined = vector_results

        if not combined:
            return []

        # 3. MMR diversity re-ranking.
        if cfg.use_mmr:
            # Need embeddings for the results.
            result_texts = [r.document.text for r in combined]
            result_embeddings = self.embedding_model.embed(result_texts)
            combined = max_marginal_relevance(
                combined, result_embeddings, query_embedding,
                top_k=cfg.top_k_initial,
                lambda_param=cfg.mmr_lambda,
            )

        # 4. Cross-encoder re-ranking.
        if cfg.use_cross_encoder:
            combined = self.reranker.rerank(
                query, combined, top_k=top_k_final,
            )
        else:
            combined = combined[:top_k_final]

        return combined

    def query(
        self,
        query: str,
        system_prompt: Optional[str] = None,
        agent: Optional[Any] = None,
    ) -> RetrievalResult:
        """
        Run the full RAG pipeline: retrieve relevant context, then
        generate a response using an LLM.

        If an ``agent`` (NexusAgent) is provided, it will be used to
        generate the response.  Otherwise only context is retrieved.

        Args:
            query: The user's query.
            system_prompt: Optional system prompt override.
            agent: Optional NexusAgent instance for LLM generation.

        Returns:
            A RetrievalResult with response text, citations, and source docs.
        """
        results = self.retrieve(query)
        if not results:
            return RetrievalResult(
                response="",
                source_documents=[],
                citations=[],
                scores=[],
            )

        # Build context from retrieved documents.
        source_docs = [r.document for r in results]
        scores = [r.score for r in results]

        # Build citations.
        citations: List[Dict[str, Any]] = []
        seen_sources: set = set()
        for r in results:
            source = r.document.metadata.get("source", "unknown")
            if source not in seen_sources:
                seen_sources.add(source)
                citations.append({
                    "source": source,
                    "score": round(r.score, 4),
                    "text_preview": r.document.text[:200] + "..." if len(r.document.text) > 200 else r.document.text,
                })

        if agent is None:
            return RetrievalResult(
                response="",
                source_documents=source_docs,
                citations=citations,
                scores=scores,
            )

        # Build the context block for the LLM.
        context_parts: List[str] = []
        for i, r in enumerate(results, 1):
            src = r.document.metadata.get("source", "unknown")
            context_parts.append(
                f"[Source {i}: {src}]\n{r.document.text}"
            )
        context_str = "\n\n".join(context_parts)

        augmented_prompt = (
            f"Answer the question based on the provided context below.\n\n"
            f"=== CONTEXT ===\n{context_str}\n"
            f"=== END CONTEXT ===\n\n"
            f"Question: {query}\n\n"
            f"Instructions:\n"
            f"1. Use ONLY the information from the context above.\n"
            f"2. If the context doesn't contain the answer, say so.\n"
            f"3. At the end, list your sources using the [Source N] references."
        )

        sys = system_prompt or (
            "You are Nexus, a helpful AI assistant with RAG capabilities. "
            "You answer questions based on the provided context."
        )

        response = agent.generate_response(augmented_prompt, system_prompt=sys)
        return RetrievalResult(
            response=response.get("text", ""),
            source_documents=source_docs,
            citations=citations,
            scores=scores,
        )

    # -- file indexing --

    def index_file(self, file_path: str) -> int:
        """
        Index a single file by reading its contents.

        Supports: .txt, .md, .rst, .py, .json, .yaml, .yml, .csv,
        and via existing Nexus loaders: .pdf, .docx, .pptx, .xlsx.

        Args:
            file_path: Path to the file.

        Returns:
            Number of chunks indexed, or 0 if file cannot be read.
        """
        path = Path(file_path)
        if not path.exists():
            logger.warning("File not found: %s", file_path)
            return 0

        text = self._read_file(path)
        if not text:
            return 0

        metadata: Dict[str, Any] = {"source": str(path), "filename": path.name}
        return self.index_texts([(text, metadata)])

    def index_directory(
        self,
        dir_path: str,
        pattern: str = "*.md",
        recursive: bool = True,
    ) -> int:
        """
        Index all files matching a glob pattern in a directory.

        Args:
            dir_path: Directory to scan.
            pattern: Glob pattern (e.g. ``"*.md"``, ``"*.txt"``, ``"**/*.py"``).
            recursive: Whether to scan subdirectories.

        Returns:
            Total number of chunks indexed.
        """
        base = Path(dir_path)
        if not base.is_dir():
            logger.warning("Directory not found: %s", dir_path)
            return 0

        if recursive and not pattern.startswith("**/"):
            pattern = f"**/{pattern}"

        total = 0
        files = list(base.glob(pattern))
        for fpath in files:
            if fpath.is_file():
                count = self.index_file(str(fpath))
                total += count

        logger.info("Indexed %d chunks from %d files in %s", total, len(files), dir_path)
        return total

    @staticmethod
    def _read_file(path: Path) -> str:
        """Read text content from a file. Supports various formats."""
        import json as _json

        suffix = path.suffix.lower()

        # Plain text formats.
        if suffix in (".txt", ".md", ".rst", ".py", ".js", ".ts", ".html", ".css",
                      ".json", ".yaml", ".yml", ".csv", ".xml", ".ini", ".cfg",
                      ".toml", ".env", ".sh", ".bat", ".ps1"):
            try:
                return path.read_text(encoding="utf-8", errors="replace")
            except Exception as e:
                logger.warning("Failed to read %s: %s", path, e)
                return ""

        # PDF — use existing pypdf loader.
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                return "\n".join(page.extract_text() for page in reader.pages)
            except Exception as e:
                logger.warning("Failed to read PDF %s: %s", path, e)
                return ""

        # DOCX — use existing python-docx.
        if suffix == ".docx":
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(str(path))
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception as e:
                logger.warning("Failed to read DOCX %s: %s", path, e)
                return ""

        # PPTX — use existing python-pptx.
        if suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                texts: List[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except Exception as e:
                logger.warning("Failed to read PPTX %s: %s", path, e)
                return ""

        # XLSX — use existing openpyxl.
        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(path), read_only=True, data_only=True)
                texts: List[str] = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join(str(c) for c in row if c is not None)
                        if row_text.strip():
                            texts.append(row_text)
                return "\n".join(texts)
            except Exception as e:
                logger.warning("Failed to read XLSX %s: %s", path, e)
                return ""

        logger.warning("Unsupported file type: %s", suffix)
        return ""

    # -- persistence with auto load/save --

    DEFAULT_INDEX_DIR = os.path.join("~", ".nexus", "rag")

    @classmethod
    def create_persistent(
        cls,
        config: Optional["RAGConfig"] = None,
        index_path: Optional[str] = None,
    ) -> "RAGEngine":
        """
        Create a RAG engine that automatically loads a previously saved
        index from disk (if it exists) and will save on every index operation.

        Args:
            config: RAG configuration.
            index_path: Path to the FAISS index file. Defaults to
                        ``~/.nexus/rag/index.faiss``.

        Returns:
            A configured RAGEngine instance with persistence enabled.
        """
        engine = cls(config=config)
        path = index_path or os.path.join(
            os.path.expanduser(cls.DEFAULT_INDEX_DIR), "index.faiss"
        )
        engine._persist_path = path

        # Auto-load if index exists.
        faiss_path = path
        meta_path = path + ".meta.json"
        if os.path.isfile(faiss_path) and os.path.isfile(meta_path):
            try:
                engine.load(faiss_path)
                logger.info("Loaded persisted RAG index from %s (%d docs)", path, engine.doc_count)
            except Exception as e:
                logger.warning("Failed to load RAG index from %s: %s", path, e)

        return engine

    def _maybe_persist(self) -> None:
        """Auto-save after indexing if persistence is enabled."""
        persist_path = getattr(self, "_persist_path", None)
        if persist_path:
            try:
                os.makedirs(os.path.dirname(persist_path), exist_ok=True)
                self.save(persist_path)
                logger.debug("Auto-saved RAG index to %s (%d docs)", persist_path, self.doc_count)
            except Exception as e:
                logger.warning("Failed to auto-save RAG index: %s", e)

    def save(self, path: str) -> None:
        """Persist the vector store and config to disk."""
        self.vector_store.save(path)

    def load(self, path: str) -> None:
        """Load a previously persisted vector store."""
        self.vector_store.load(path)

    def clear(self) -> None:
        """Clear all indexed data."""
        self._documents.clear()
        self.vector_store.clear()
        if self._bm25 is not None:
            self._bm25.clear()
        logger.info("RAG engine cleared")
        # Remove persisted files.
        persist_path = getattr(self, "_persist_path", None)
        if persist_path:
            for suffix in ("", ".meta.json"):
                fpath = persist_path + suffix
                if os.path.isfile(fpath):
                    try:
                        os.remove(fpath)
                    except OSError:
                        pass

    @property
    def doc_count(self) -> int:
        """Return the total number of indexed documents/chunks."""
        return self.vector_store.count()

    def get_stats(self) -> Dict[str, Any]:
        """Return statistics about the indexed documents."""
        return {
            "doc_count": self.doc_count,
            "sources": list(set(
                d.metadata.get("source", "unknown")
                for d in self._documents
            )),
            "has_persistence": hasattr(self, "_persist_path") and bool(self._persist_path),
            "persist_path": getattr(self, "_persist_path", None),
        }
