"""
Content Loader Module

Provides functions to load content from various sources:
webpages, YouTube videos, PDF, DOCX, PPTX, Excel, and plain text files.
"""

import io
import logging
import os
import re
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------


def _fetch_url(url: str, timeout: int = 30) -> str:
    """Download a URL and return its raw text content."""
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/120.0.0.0 Safari/537.36"
        )
    }
    resp = requests.get(url, headers=headers, timeout=timeout)
    resp.raise_for_status()
    return resp.text


def _download_binary(url: str, timeout: int = 30) -> bytes:
    """Download a binary file from a URL."""
    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/120.0.0.0 Safari/537.36"
        )
    }
    resp = requests.get(url, headers=headers, timeout=timeout)
    resp.raise_for_status()
    return resp.content


def _is_url(path: str) -> bool:
    """Return True if *path* looks like a URL."""
    parsed = urlparse(path)
    return parsed.scheme in ("http", "https")


# ---------------------------------------------------------------------------
# Individual loaders
# ---------------------------------------------------------------------------


def load_webpage(url: str) -> str:
    """Load text content from a webpage using requests + BeautifulSoup."""
    logger.info("Loading webpage: %s", url)
    try:
        html = _fetch_url(url)
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        content = "\n".join(lines)
        logger.debug("Webpage loaded: %d characters", len(content))
        return content
    except requests.RequestException as e:
        logger.error("Failed to load webpage %s: %s", url, e)
        return f"[Ошибка загрузки веб-страницы: {e}]"


def load_youtube(url: str) -> str:
    """Load transcript from a YouTube video."""
    logger.info("Loading YouTube video: %s", url)
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        video_id = None
        patterns = [
            r"v=([a-zA-Z0-9_-]{11})",
            r"youtu\.be/([a-zA-Z0-9_-]{11})",
            r"embed/([a-zA-Z0-9_-]{11})",
        ]
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                video_id = match.group(1)
                break
        if not video_id:
            return "[Ошибка: не удалось извлечь ID видео из URL]"
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        transcript_text = " ".join(entry["text"] for entry in transcript_list)
        logger.debug("YouTube transcript loaded: %d characters", len(transcript_text))
        return f"YouTube video transcript:\n\n{transcript_text}"
    except ImportError:
        logger.error("youtube-transcript-api is not installed")
        return "[Ошибка: youtube-transcript-api не установлен]"
    except Exception as e:
        logger.error("Failed to load YouTube video %s: %s", url, e)
        return f"[Ошибка загрузки YouTube видео: {e}]"


def load_pdf(url_or_path: str) -> str:
    """Load text from a PDF file (local path or URL)."""
    logger.info("Loading PDF: %s", url_or_path)
    try:
        from pypdf import PdfReader
        if _is_url(url_or_path):
            raw = _download_binary(url_or_path)
            reader = PdfReader(io.BytesIO(raw))
        else:
            reader = PdfReader(url_or_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append(f"--- Страница {i + 1} ---\n{text.strip()}")
        content = "\n\n".join(pages)
        logger.debug("PDF loaded: %d pages, %d characters", len(pages), len(content))
        return content
    except ImportError:
        logger.error("pypdf is not installed")
        return "[Ошибка: pypdf не установлен]"
    except Exception as e:
        logger.error("Failed to load PDF %s: %s", url_or_path, e)
        return f"[Ошибка загрузки PDF: {e}]"


def load_docx(url_or_path: str) -> str:
    """Load text from a DOCX file (local path or URL)."""
    logger.info("Loading DOCX: %s", url_or_path)
    try:
        from docx import Document
        if _is_url(url_or_path):
            raw = _download_binary(url_or_path)
            doc = Document(io.BytesIO(raw))
        else:
            doc = Document(url_or_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n".join(paragraphs)
        logger.debug("DOCX loaded: %d characters", len(content))
        return content
    except ImportError:
        logger.error("python-docx is not installed")
        return "[Ошибка: python-docx не установлен]"
    except Exception as e:
        logger.error("Failed to load DOCX %s: %s", url_or_path, e)
        return f"[Ошибка загрузки DOCX: {e}]"


def load_pptx(url_or_path: str) -> str:
    """Load text from a PPTX file (local path or URL)."""
    logger.info("Loading PPTX: %s", url_or_path)
    try:
        from pptx import Presentation
        if _is_url(url_or_path):
            raw = _download_binary(url_or_path)
            prs = Presentation(io.BytesIO(raw))
        else:
            prs = Presentation(url_or_path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_lines = [f"--- Слайд {i + 1} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_lines.append(text)
            slides_text.append("\n".join(slide_lines))
        content = "\n\n".join(slides_text)
        logger.debug("PPTX loaded: %d slides, %d characters", len(prs.slides), len(content))
        return content
    except ImportError:
        logger.error("python-pptx is not installed")
        return "[Ошибка: python-pptx не установлен]"
    except Exception as e:
        logger.error("Failed to load PPTX %s: %s", url_or_path, e)
        return f"[Ошибка загрузки PPTX: {e}]"


def load_excel(url_or_path: str) -> str:
    """Load text from an Excel file (local path or URL) using openpyxl."""
    logger.info("Loading Excel: %s", url_or_path)
    try:
        import openpyxl
        if _is_url(url_or_path):
            raw = _download_binary(url_or_path)
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        else:
            wb = openpyxl.load_workbook(url_or_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows():
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                rows.append("\t".join(row_values))
            sheet_text = f"--- Лист: {sheet_name} ---\n" + "\n".join(rows)
            parts.append(sheet_text)
        content = "\n\n".join(parts)
        wb.close()
        logger.debug("Excel loaded: %d sheets, %d characters", len(wb.sheetnames), len(content))
        return content
    except ImportError:
        logger.error("openpyxl is not installed")
        return "[Ошибка: openpyxl не установлен]"
    except Exception as e:
        logger.error("Failed to load Excel %s: %s", url_or_path, e)
        return f"[Ошибка загрузки Excel: {e}]"


def load_text_file(url_or_path: str) -> str:
    """Load a plain text file (local path or URL)."""
    logger.info("Loading text file: %s", url_or_path)
    try:
        if _is_url(url_or_path):
            return _fetch_url(url_or_path)
        else:
            with open(url_or_path, "r", encoding="utf-8") as fh:
                return fh.read()
    except Exception as e:
        logger.error("Failed to load text file %s: %s", url_or_path, e)
        return f"[Ошибка загрузки текстового файла: {e}]"


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_EXTENSION_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_excel,
    ".xls": load_excel,
    ".txt": load_text_file,
    ".md": load_text_file,
    ".csv": load_text_file,
    ".json": load_text_file,
    ".xml": load_text_file,
    ".yaml": load_text_file,
    ".yml": load_text_file,
    ".html": load_webpage,
    ".htm": load_webpage,
}


def load(url: str) -> str:
    """
    Detect the content type by URL / file extension and load the content.

    Supported: web pages, YouTube, PDF, DOCX, PPTX, Excel, text files.
    """
    logger.info("Loading content from: %s", url)

    # YouTube detection
    youtube_patterns = [
        r"youtube\.com/watch\?v=",
        r"youtu\.be/",
        r"youtube\.com/embed/",
        r"youtube\.com/shorts/",
    ]
    for pattern in youtube_patterns:
        if re.search(pattern, url):
            return load_youtube(url)

    # Extension-based dispatch
    lower = url.lower()
    for ext, loader_func in _EXTENSION_MAP.items():
        if lower.endswith(ext):
            return loader_func(url)

    # Default: treat as webpage
    if _is_url(url):
        return load_webpage(url)

    # Last resort: local file
    if os.path.isfile(url):
        return load_text_file(url)

    logger.warning("Unknown content type for: %s", url)
    return f"[Неизвестный тип контента: {url}]"